import sys, os
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

C_DARK = RGBColor(0x0a, 0x2a, 0x44)
C_MID = RGBColor(0x0a, 0x3d, 0x62)
C_LIGHT = RGBColor(0x1a, 0x6b, 0x8a)
C_ACCENT = RGBColor(0x4f, 0xcd, 0xa5)
C_GOLD = RGBColor(0xff, 0xd7, 0x00)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_OFF = RGBColor(0xE8, 0xF4, 0xF8)
C_ORANGE = RGBColor(0xFF, 0x8c, 0x42)
C_RED = RGBColor(0xef, 0x53, 0x50)
C_GREEN = RGBColor(0x66, 0xBB, 0x6a)
C_PURPLE = RGBColor(0xab, 0x47, 0xbc)

def add_bg(slide, top=C_MID, bottom=C_LIGHT):
    w, h = prs.slide_width, prs.slide_height
    for i in range(15):
        r = top[0] + int((bottom[0] - top[0]) * i / 14)
        g = top[1] + int((bottom[1] - top[1]) * i / 14)
        b = top[2] + int((bottom[2] - top[2]) * i / 14)
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, int(h * i / 15), w, int(h / 15) + 1)
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(r, g, b)
        rect.line.fill.background()

def add_box(slide, l, t, w, h, color, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def add_txt(slide, l, t, w, h, text, size=20, color=C_WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tb

def add_lines(slide, l, t, w, h, lines, size=20, color=C_OFF, gap=Pt(6)):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = gap
    return tb

def add_title_bar(slide, title, subtitle=None):
    add_box(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), C_ACCENT)
    add_txt(slide, Inches(0.8), Inches(0.1), Inches(11), Inches(0.7), title, size=34, color=C_DARK, bold=True)
    if subtitle:
        add_txt(slide, Inches(0.8), Inches(0.65), Inches(11), Inches(0.4), subtitle, size=16, color=C_DARK)

def add_section_label(slide, text, y=Inches(1.3)):
    add_box(slide, Inches(0), y, Inches(0.15), Inches(0.5), C_ACCENT)
    add_txt(slide, Inches(0.3), y, Inches(3), Inches(0.5), text, size=18, color=C_ACCENT, bold=True)

def add_quiz_box(slide, title, scenario, options, answer, reason, errors, y_start=Inches(1.8)):
    add_box(slide, Inches(0.5), y_start, Inches(12.3), Inches(5.2), C_MID, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_txt(slide, Inches(0.8), y_start + Inches(0.15), Inches(11.5), Inches(0.5), title, size=24, color=C_GOLD, bold=True)
    
    current_y = y_start + Inches(0.65)
    add_txt(slide, Inches(0.8), current_y, Inches(11.5), Inches(0.8), scenario, size=20, color=C_WHITE)
    current_y += Inches(0.9)
    
    for opt in options:
        add_txt(slide, Inches(1.0), current_y, Inches(11), Inches(0.4), opt, size=18, color=C_OFF)
        current_y += Inches(0.38)
    
    current_y += Inches(0.15)
    add_box(slide, Inches(0.8), current_y, Inches(11.5), Inches(0.04), C_GOLD)
    current_y += Inches(0.15)
    
    add_txt(slide, Inches(0.8), current_y, Inches(11.5), Inches(0.35), f"正解：{answer}", size=20, color=C_GREEN, bold=True)
    current_y += Inches(0.35)
    add_txt(slide, Inches(0.8), current_y, Inches(11.5), Inches(0.7), f"理由：{reason}", size=17, color=C_OFF)
    current_y += Inches(0.5)
    
    add_box(slide, Inches(0.8), current_y, Inches(0.08), Inches(0.6), C_RED)
    add_txt(slide, Inches(1.0), current_y, Inches(11), Inches(0.6), f"常見錯誤：{errors}", size=17, color=C_ORANGE)

def add_footer(slide, num, total=10):
    add_txt(slide, Inches(12), Inches(7.0), Inches(1), Inches(0.4), f"{num}/{total}", size=12, color=C_OFF)

# ===================== P1: 導入 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, C_DARK, C_MID)
add_txt(s, Inches(1), Inches(0.8), Inches(11.3), Inches(0.5), "公民與社會 | 高三複習", size=18, color=C_ACCENT, align=PP_ALIGN.CENTER)
add_txt(s, Inches(1), Inches(1.6), Inches(11.3), Inches(1.0), "權力分立與制衡\n——情境判斷這樣學", size=44, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_box(s, Inches(4.5), Inches(3.2), Inches(4.3), Inches(0.06), C_GOLD)
add_txt(s, Inches(1), Inches(3.6), Inches(11.3), Inches(0.8), "學習路徑：為什麼要分權 → 三權怎麼互動 → 情境題怎麼解", size=22, color=C_OFF, align=PP_ALIGN.CENTER)
add_lines(s, Inches(1.5), Inches(4.8), Inches(10.3), Inches(2), [
    "本節目標：",
    "  \u25b8 理解權力分立的目的（不是為了吵架，是為了避免獨裁）",
    "  \u25b8 能判斷真實情境中哪個權力在運作（行政？立法？司法？）",
    "  \u25b8 能辨識「制衡」關係與「權力越界」的常見考題陷阱",
], size=20, color=C_OFF)
add_footer(s, 1)

# ===================== P2: 迷思破除 + 互動 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title_bar(s, "權力分立不是「三權各自為政」！", "先來破除最大迷思")
add_lines(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(1.5), [
    "核心觀念：",
    "  \u25b8 權力分立（Separation of Powers）：行政、立法、司法各自獨立",
    "  \u25b8 權力制衡（Checks and Balances）：彼此監督、互相節制",
    "  \u25b8 目的：防止政府權力過度集中 → 保障人民權利",
], size=22, color=C_OFF)
add_quiz_box(s, "\U0001f3b2 情境選擇題 \u2502 第1題：迷思檢測", 
    "新聞報導：「行政院長在立法院備詢時，立法委員對其施政提出質疑。」\n這屬於哪種權力互動？",
    ["A. 行政權干預立法權", "B. 立法權監督行政權", "C. 司法權介入行政權", "D. 三權完全獨立互不往來"],
    "B",
    "立法委員質詢行政院長，是立法院行使「監督權」的具體表現，屬於立法權對行政權的制衡。\nA：行政院長是被監督方，不是干預方。C：這裡沒有法院或法官。D：實際上權力間有互動，不是完全不往來。",
    "同學常誤以為「質詢=行政與立法吵架」，但這是憲法賦予立法院的合法監督工具，是正常制衡運作。不要看到衝突就直覺選「不應該」。",
    y_start=Inches(3.0))
add_footer(s, 2)

# ===================== P3: 三權職掌表 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title_bar(s, "三權分立的職掌與角色", "先搞清楚誰在做什麼，才能判斷有沒有越界")
add_box(s, Inches(0.8), Inches(1.5), Inches(3.5), Inches(2.5), C_ORANGE, MSO_SHAPE.ROUNDED_RECTANGLE)
add_txt(s, Inches(0.8), Inches(1.6), Inches(3.5), Inches(0.4), "\U0001f3e2 行政權", size=24, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_txt(s, Inches(0.8), Inches(2.1), Inches(3.5), Inches(1.5), "執行法律、推動政策\n代表：總統、行政院\n口訣：做事的人", size=18, color=C_OFF, align=PP_ALIGN.CENTER)

add_box(s, Inches(4.8), Inches(1.5), Inches(3.5), Inches(2.5), C_PURPLE, MSO_SHAPE.ROUNDED_RECTANGLE)
add_txt(s, Inches(4.8), Inches(1.6), Inches(3.5), Inches(0.4), "\U0001f4da 立法權", size=24, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_txt(s, Inches(4.8), Inches(2.1), Inches(3.5), Inches(1.5), "制定法律、審查預算\n代表：立法院\n口訣：訂規則的人", size=18, color=C_OFF, align=PP_ALIGN.CENTER)

add_box(s, Inches(8.8), Inches(1.5), Inches(3.5), Inches(2.5), C_GREEN, MSO_SHAPE.ROUNDED_RECTANGLE)
add_txt(s, Inches(8.8), Inches(1.6), Inches(3.5), Inches(0.4), "\u2696\ufe0f 司法權", size=24, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_txt(s, Inches(8.8), Inches(2.1), Inches(3.5), Inches(1.5), "解釋法律、審判案件\n代表：司法院、法官\n口訣：評理的人", size=18, color=C_OFF, align=PP_ALIGN.CENTER)

add_box(s, Inches(0.8), Inches(4.3), Inches(11.5), Inches(0.04), C_GOLD)
add_lines(s, Inches(0.8), Inches(4.6), Inches(11.5), Inches(1.5), [
    "制衡關係速記：行政做事 → 立法審查（預算+法律）→ 司法裁判（是否違法違憲）",
    "行政可以提出法律案（提案權），但需要立法院通過；立法院通過後，行政必須執行。",
], size=18, color=C_OFF)
add_quiz_box(s, "\U0001f3b2 第2題：職掌判斷",
    "某縣市政府公告「禁止餐廳提供一次性塑膠吸管」，違者罰款。這屬於哪一權的行使？",
    ["A. 立法權（因為涉及罰款規定）", "B. 司法權（因為要執行處罰）", "C. 行政權（因為是政府執行法律）", "D. 監察權（因為要監督店家）"],
    "C",
    "縣市政府發布行政命令並執行處罰，是行政機關執行法律的表現。\n雖然涉及罰款，但這是「依法行政」，並非制定新法律（那是立法院的事）。",
    "看到「罰款」就以為是司法權是大忌！行政機關也可以依法處罰（例如環保局開罰單），司法權是「事後審判」，不是開罰單。",
    y_start=Inches(5.0))
add_footer(s, 3)

# ===================== P4: 制衡機制 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title_bar(s, "權力制衡怎麼運作？", "三個權力之間的「你中有我、我中有你」")
add_lines(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(2.5), [
    "行政 \u2194 立法：",
    "  \u25b8 立法審查預算（砍預算就是制衡）",
    "  \u25b8 立法質詢官員（要求來說明）",
    "  \u25b8 行政可提法律案、可覆議（退回法律請立院重審）",
    "",
    "行政 \u2194 司法：",
    "  \u25b8 司法審查行政處分是否違法（行政訴訟）",
    "",
    "立法 \u2194 司法：",
    "  \u25b8 大法官解釋法律是否違憲（司法院釋憲）",
], size=20, color=C_OFF)
add_quiz_box(s, "\U0001f3b2 第3題：制衡判斷",
    "立法院通過了《礦業法》修正案，但行政院認為部分條文窒礙難行，於是提出覆議。\n請問這是什麼權力的運作？",
    ["A. 行政權侵犯立法權", "B. 立法權監督行政權", "C. 行政權對立法權的制衡", "D. 司法權介入立法程序"],
    "C",
    "行政院對立法院通過的法律提出覆議，是憲法賦予行政權的制衡工具，要求立法院重新審議。\n這不是侵犯，而是合法制衡。A的「侵犯」是錯誤用詞，因為這是憲法允許的。",
    "看到行政院「反對」立法院，容易直覺選「侵犯立法權」。但覆議是憲法明文規定的制衡機制，不是違法干預。",
    y_start=Inches(4.0))
add_footer(s, 4)

# ===================== P5: 互動題1 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title_bar(s, "類題練習 ①：情境判斷", "老師引導思考——先找關鍵字，再對應權力")
add_lines(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(1.2), [
    "解題三步驟：",
    "  Step 1：看主詞（誰在做？——總統？立委？法官？）",
    "  Step 2：看動作（做什麼？——開罰？修法？判決？）",
    "  Step 3：對表格（屬於哪一權？有沒有制衡關係？）",
], size=22, color=C_GOLD)
add_quiz_box(s, "\U0001f3b2 第4題",
    "某直轄市長認為議會通過的地方自治條例有違憲疑義，\n向司法院聲請解釋憲法。請問這屬於？",
    ["A. 行政權監督立法權", "B. 司法權被動啟動釋憲", "C. 立法權審查行政權", "D. 監察權調查地方自治"],
    "B",
    "市長（行政）認為法規違憲，向司法院大法官聲請釋憲，是啟動司法權的違憲審查機制。\n司法權是被動的——不告不理，必須有人聲請才能啟動。",
    "易誤選A，因為看起來是行政對立法有意見。但這裡的關鍵是「聲請釋憲」，動作是交給司法機關裁判，不是行政自己否定法律。",
    y_start=Inches(2.6))
add_footer(s, 5)

# ===================== P6: 互動題2 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title_bar(s, "類題練習 ②：越界或合法？", "區分「合法制衡」與「權力越界」")
add_lines(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(1.0), [
    "核心判斷：合法制衡是憲法允許的；越界是把別人的權力抢來自己做。",
    "  \u25b8 合法制衡：質詢、審預算、覆議、釋憲",
    "  \u25b8 權力越界：行政自己宣布法律無效、立法代替法院判決",
], size=20, color=C_GOLD)
add_quiz_box(s, "\U0001f3b2 第5題",
    "某立法委員在質詢時要求行政院長「立刻撤換某位部長」，\n否則將刪減該部會預算。此舉是否恰當？",
    ["A. 完全恰當，立法權本來就可以干預人事", "B. 不恰當，人事任命權屬於行政權，立法不能直接強迫", "C. 恰當，因為預算權是立法院的權力", "D. 不恰當，應由司法權決定人事"],
    "B",
    "立法委員可以質詢、可以審預算，但不能直接命令行政院長撤換部長（人事任命權屬行政權）。\n用預算威脅要求撤換，已經逾越監督範圍，侵犯行政權的人事決定權。",
    "易誤選C——以為「有預算權就可以為所欲為」。預算權是審查政府花費，不是拿來交換人事的手段。",
    y_start=Inches(2.5))
add_footer(s, 6)

# ===================== P7: 互動題3 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title_bar(s, "類題練習 ③：司法權的特殊性", "司法權是「被動」、「中立」、「個案」")
add_lines(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(1.2), [
    "司法權三特性（考試最愛考）：",
    "  \u25b8 被動性：不告不理（沒有人起訴，法院不能主動審判）",
    "  \u25b8 中立性：法官獨立審判，不受行政或立法干涉",
    "  \u25b8 個案性：只處理具體案件，不主動解釋抽象法律問題",
], size=20, color=C_GOLD)
add_quiz_box(s, "\U0001f3b2 第6題",
    "某民間團體認為《道路交通管理處罰條例》某條文違憲，\n直接向大法官聲請釋憲，但被駁回。為什麼？",
    ["A. 大法官不想處理這個案子", "B. 民間團體不能直接聲請釋憲，須經由法院", "C. 條文根本沒有違憲", "D. 大法官只處理行政機關的聲請"],
    "B",
    "根據現行制度，人民不能直接聲請釋憲，必須在訴訟過程中由「承審法院」認為法律有違憲疑義，\n裁定停止訴訟後，才能由法院向大法官聲請解釋。這就是司法權的被動性。",
    "常有人以為「我覺得違憲就可以直接找大法官」，但大法官不受理個人的直接聲請（有例外但很少）。\n關鍵：司法權必須透過「案件」才能啟動。",
    y_start=Inches(2.6))
add_footer(s, 7)

# ===================== P8: 互動題4 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title_bar(s, "跨章節整合：我國中央政府體制", "總統 vs 行政院長，權力怎麼分？")
add_lines(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(1.8), [
    "我國（雙首長制）特色：",
    "  \u25b8 總統：國家元首，國防外交，任命行政院長（不需立法院同意）",
    "  \u25b8 行政院長：最高行政首長，對立法院負責",
    "  \u25b8 立法院：可對行政院長提出不信任案（倒閣）",
    "  \u25b8 總統：在立法院倒閣後，可解散立法院（反制）",
    "",
    "這整套機制就是「行政 vs 立法」的完整制衡設計！",
], size=20, color=C_OFF)
add_quiz_box(s, "\U0001f3b2 第7題：跨章節應用",
    "立法院通過對行政院長的不信任案後，\n行政院長應如何回應？（依我國憲法增修條文）",
    ["A. 行政院長必須辭職", "B. 總統可以解散立法院", "C. 行政院長可以提覆議", "D. A和B都是正確的"],
    "D",
    "不信任案通過後，行政院長須辭職（A）；總統可在十天內諮詢立法院長後，宣告解散立法院（B）。\n這是憲法設計的制衡：立法院可以倒閣，總統可以解散國會。",
    "很多人只記得「倒閣=行政院長下台」，忘了總統還有解散立法院的反制權力。制衡是雙向的！",
    y_start=Inches(3.3))
add_footer(s, 8)

# ===================== P9: 常見錯誤總整理 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title_bar(s, "常見錯誤總整理", "")

add_box(s, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.0), C_MID, MSO_SHAPE.ROUNDED_RECTANGLE)
add_txt(s, Inches(0.8), Inches(1.6), Inches(5.2), Inches(0.4), "\u274c 錯誤觀念", size=22, color=C_RED, bold=True, align=PP_ALIGN.CENTER)
add_lines(s, Inches(0.8), Inches(2.1), Inches(5.2), Inches(4.0), [
    "1. 質詢 = 行政立法吵架（X）",
    "   \u2192 質詢是合法監督機制",
    "",
    "2. 有罰款就是司法權（X）",
    "   \u2192 行政機關也可以依法行政開罰",
    "",
    "3. 預算權可以威脅人事（X）",
    "   \u2192 預算審查≠人事命令權",
    "",
    "4. 大法官可以直接受理人民（X）",
    "   \u2192 需經法院審理中才能聲請",
], size=17, color=C_OFF)

add_box(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.0), C_MID, MSO_SHAPE.ROUNDED_RECTANGLE)
add_txt(s, Inches(7.1), Inches(1.6), Inches(5.2), Inches(0.4), "\u2705 正確觀念", size=22, color=C_GREEN, bold=True, align=PP_ALIGN.CENTER)
add_lines(s, Inches(7.1), Inches(2.1), Inches(5.2), Inches(4.0), [
    "1. 質詢是立法權制衡行政權",
    "",
    "2. 看主詞：行政機關開罰＝行政權",
    "   法院判決＝司法權",
    "",
    "3. 人事任命權專屬行政權",
    "   立法只能用法律規範資格",
    "",
    "4. 司法權被動，不告不理",
    "   須透過案件程序才能啟動",
], size=17, color=C_OFF)

add_footer(s, 9)

# ===================== P10: 總結 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, C_DARK, C_MID)
add_txt(s, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.5), "本課總結：權力分立與制衡 答題思考流程", size=28, color=C_ACCENT, bold=True, align=PP_ALIGN.CENTER)

add_box(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(3.2), C_MID, MSO_SHAPE.ROUNDED_RECTANGLE)
add_lines(s, Inches(0.8), Inches(1.4), Inches(11.7), Inches(2.8), [
    "遇到情境題，三步驟解題：",
    "",
    "\u2460 找主詞：是誰在做？（總統/市長/立委/法官/大法官）",
    "\u2461 找動作：做什麼？（開罰/修法/質詢/釋憲/判決/覆議）",
    "\u2462 對權力表：行政做事+執行／立法制定+審查／司法審判+解釋",
    "",
    "進階判斷：是「合法制衡」還是「權力越界」？",
    "  \u2192 質詢、審預算、覆議、釋憲、不信任案、解散國會 → 合法",
    "  \u2192 行政拒執行法律、立法強迫人事、法院主動調查 → 越界",
], size=20, color=C_OFF)

add_quiz_box(s, "\U0001f3b2 最後一題：總檢核",
    "某法官在審理一起酒駕案件時，認為《刑法》某條文違憲，\n於是裁定停止審判，向大法官聲請釋憲。請問這說明了什麼？",
    ["A. 司法權的主動性", "B. 立法權優於司法權", "C. 司法權的違憲審查啟動方式", "D. 行政權干預司法審判"],
    "C",
    "法官在審理具體案件時，若認為適用法律有違憲疑義，可裁定停止審判並聲請釋憲。\n這正是司法權違憲審查的啟動方式，也體現了司法權的被動性（須有案件在先）。\nA錯在司法權不是主動的；B錯在法律是否違憲由大法官判斷，不是立法權說了算。",
    "法官聲請釋憲常被誤以為是「法官不認同法律」，其實這是憲法賦予法官的義務——\n當發現法律可能違憲時，法官有責任暫停審判並聲請解釋。",
    y_start=Inches(4.6))
add_footer(s, 10)

# Save
output = os.path.expanduser(r"~\Documents\emoji-collection\高三公民_權力分立與制衡_教學PPT.pptx")
prs.save(output)
print(f"PPT saved: {output}")
print(f"Size: {os.path.getsize(output) / 1024:.0f} KB")
print(f"Slides: {len(prs.slides)}")
