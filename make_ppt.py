import sys, os
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BG_TOP = RGBColor(0x0a, 0x3d, 0x62)
BG_BOTTOM = RGBColor(0x1a, 0x6b, 0x8a)
ACCENT = RGBColor(0x4f, 0xcd, 0xa5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xE8, 0xF4, 0xF8)
DARK = RGBColor(0x0a, 0x2a, 0x44)
GOLD = RGBColor(0xFF, 0xD7, 0x00)

def add_bg(slide, top=BG_TOP, bottom=BG_BOTTOM):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    w, h = prs.slide_width, prs.slide_height
    for i in range(20):
        r = top[0] + int((bottom[0] - top[0]) * i / 19)
        g = top[1] + int((bottom[1] - top[1]) * i / 19)
        b = top[2] + int((bottom[2] - top[2]) * i / 19)
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, int(h * i / 20), w, int(h / 20) + 1
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(r, g, b)
        rect.line.fill.background()

def add_deco(slide):
    deco = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(-1), Inches(4), Inches(4))
    deco.fill.solid()
    deco.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    deco.fill.fore_color.brightness = 0.95
    deco.line.fill.background()

def add_title_text(slide, text, size=40, y=Inches(2.8), color=WHITE, bold=True):
    txBox = slide.shapes.add_textbox(Inches(1), y, Inches(11.3), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = PP_ALIGN.CENTER
    return txBox

def add_subtitle(slide, text, y=Inches(4.2), size=20):
    txBox = slide.shapes.add_textbox(Inches(2), y, Inches(9.3), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = LIGHT
    p.alignment = PP_ALIGN.CENTER
    return txBox

def add_body(slide, items, y=Inches(3.2), size=22):
    txBox = slide.shapes.add_textbox(Inches(1.2), y, Inches(10.9), Inches(4))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = LIGHT
        p.space_after = Pt(12)
        p.alignment = PP_ALIGN.LEFT
    return txBox

def add_section_title(slide, section, title):
    add_bg(slide, BG_BOTTOM, BG_TOP)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.3), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = section
    p.font.size = Pt(18)
    p.font.color.rgb = ACCENT
    p.font.bold = True
    add_title_text(slide, title, size=44, y=Inches(2.5))

def add_content_slide(slide, title, body_items, size=22):
    add_bg(slide)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.3), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.color.rgb = ACCENT
    p.font.bold = True
    add_body(slide, body_items, y=Inches(1.5), size=size)

# ========== SLIDE 1: Title ==========
s1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s1, BG_TOP, BG_BOTTOM)
add_deco(s1)

# Globe decoration
globe = s1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.5), Inches(1), Inches(2.3), Inches(2.3))
globe.fill.solid()
globe.fill.fore_color.rgb = ACCENT
globe.line.fill.background()
globe.fill.fore_color.brightness = 0.3

# Title
txBox = s1.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.3), Inches(0.6))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "給年輕世代的永續素養課"
p.font.size = Pt(22)
p.font.color.rgb = ACCENT
p.alignment = PP_ALIGN.CENTER

add_title_text(s1, "COP26 氣候峰會\n與我們想要的未來", size=48, y=Inches(2.0))
add_subtitle(s1, "永續發展 × 氣候行動 × 青年參與", y=Inches(4.5))
add_subtitle(s1, "對象：高三學生　　　　講師：陳惠萍", y=Inches(5.3), size=18)

# ========== SLIDE 2: 課程架構 ==========
s2 = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(s2, "課程架構", "本課程將帶你探索")
add_body(s2, [
    "🌍 一、融冰遊戲體驗 — 暖化與家園的距離",
    "🌡️ 二、全球暖化介紹 — 地球為什麼發燒？",
    "🏛️ 三、全球氣候峰會 — COP26 做了什麼？",
    "⚖️ 四、氣候正義行動 — 沒有人是局外人",
    "🌱 五、青年行動 — 你的選擇決定未來",
], y=Inches(3.5), size=26)

# ========== SLIDE 3: 融冰遊戲 ==========
s3 = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(s3, "🌊 融冰體驗：暖化與北極熊的家園", [
    "▸ 全班分組（4人一組），站在報紙（冰塊）上",
    "▸ 每次升溫 0.5°C → 報紙對折一次（冰層融化）",
    "▸ 升溫 1.0°C → 再對折（面積僅剩 1/4）",
    "▸ 升溫 1.5°C → 再對折（面積僅剩 1/8）",
    "▸ 升溫 2.0°C → 再對折（誰還站得住？）",
    "",
    "💭 反思：如果你是北極熊，你希望人類怎麼做？",
])

# ========== SLIDE 4: 全球暖化 ==========
s4 = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(s4, "🌡️ 地球為什麼發燒？", [
    "▸ 人類燃燒化石燃料（汽車、飛機、工廠）→ 排放 CO₂",
    "▸ 溫室氣體困住太陽熱能 → 地球逐步升溫",
    "▸ 全球平均溫度已比 1850 年代明顯上升",
    "▸ 後果：冰層融化 → 海平面上升 → 極端氣候",
    "",
    "🔑 關鍵數字：1.5°C",
    "   這是科學家認為的「安全升溫上限」",
])

# ========== SLIDE 5: 1.5°C 的重要性 ==========
s5 = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(s5, "🎯 為什麼是 1.5°C？", [
    "▸ 極端溫度：相較 2°C，1.5°C 可減少 17 億人受影響",
    "▸ 海平面上升：減少 10 公分（影響減少 1000 萬人）",
    "▸ 糧食危機：升溫使主要作物產量下降",
    "▸ 珊瑚礁：升溫 1.5°C → 70-90% 退化；2°C → 99%",
    "▸ 經濟：1.5°C 路徑可降低 25% 的全球損失成本",
    "",
    "⚡ 目標：2030 年碳排減半，2050 年達到淨零",
])

# ========== SLIDE 6: COP26 ==========
s6 = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(s6, "🏛️ COP26：拯救地球的最後機會", [
    "▸ COP = Conference of the Parties（締約國會議）",
    "▸ 1995 年起每年舉辦，197 個國家參與",
    "▸ COP26（2021，英國格拉斯哥）被視為關鍵轉折點",
    "",
    "📊 與會規模：",
    "   • 25,000 名來自 200 國的代表",
    "   • 10,000 名英國警察",
    "   • 100,000 名抗議民眾",
])

# ========== SLIDE 7: COP26 四大目標 ==========
s7 = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(s7, "🎯 COP26 四大核心目標", [
    "① 控制全球升溫不超過 1.5°C",
    "   → 各國提出更積極的減碳承諾（NDC）",
    "",
    "② 提供 1,000 億美元氣候資金",
    "   → 已開發國家資助開發中國家轉型",
    "",
    "③ 保護生態系統與自然棲息地",
    "   → 森林保育、生物多樣性",
    "",
    "④ 政府、企業與民間的合作",
    "   → 全民參與才有真正的改變",
])

# ========== SLIDE 8: 淨零排放 ==========
s8 = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(s8, "🌱 什麼是「淨零排放」？", [
    "▸ 淨零（Net Zero）≠ 零排放",
    "▸ 定義：人為排放量 = 人為移除量，達成平衡",
    "▸ 排放來源：CO₂、甲烷、N₂O 等七種溫室氣體",
    "▸ 移除方式：森林碳匯、碳捕集技術（CCUS）",
    "",
    "📈 全球響應：",
    "   • 128+ 國家承諾淨零",
    "   • 235+ 城市",
    "   • 699+ 企業",
    "   • 台灣也宣示 2050 淨零轉型",
])

# ========== SLIDE 9: 能源轉型 ==========
s9 = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(s9, "⚡ 淨零之戰：能源轉型是關鍵", [
    "▸ 能源部門占全球碳排最大宗",
    "▸ 化石燃料（煤、石油、天然氣）→ 再生能源",
    "",
    "☀️ 全球超過 100 個城市電力 70% 來自再生能源",
    "",
    "🏙️ 格拉斯哥（Glasgow）— 綠色城市的典範",
    "   • 工業革命的搖籃 → 第四波綠色工業革命",
    "   • 蘇格蘭目標 2045 淨零",
    "   • 2021 年 95% 電力來自再生能源",
    "   • 2020 年獲選 Global Green City",
])

# ========== SLIDE 10: 氣候正義 ==========
s10 = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(s10, "⚖️ 氣候正義：誰承受最多？", [
    "🌏 吐瓦魯 — 第一個面臨消失的國家",
    "   • 平均海拔僅 198 公分，水位每年上升 0.5 公分",
    "   • 外交部長站在海水中演講，震撼全球",
    "",
    "👩 性別與氣候",
    "   • 氣候災民中 80% 是女性",
    "   • 低收入國家至少 400 萬女孩因氣候事件失學",
    "",
    "💡 核心問題：",
    "   造成暖化最少的國家，承受最大的苦難",
])

# ========== SLIDE 11: 青年力量 ==========
s11 = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(s11, "✊ 年輕世代站出來：No More Blah Blah", [
    "🌍 Greta Thunberg（格蕾塔·童貝里）",
    "   • 2018 年，15 歲，在瑞典議會前「為氣候罷課」",
    "   • 啟發全球 400 萬人參與氣候遊行",
    "   • 2019 年《時代》雜誌風雲人物（史上最年輕）",
    "",
    "🗣️ 在 COP26 上說出心聲：",
    "   「COP 已變成公關活動，領導人只會說漂亮話。」",
    '   — "Blah, blah, blah" 演講成為經典',
    "",
    "💪 青年的力量（Power of Youth）正在改變世界",
])

# ========== SLIDE 12: 台灣的角色 ==========
s12 = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(s12, "🇹🇼 台灣的氣候挑戰與行動", [
    "⚠️ 台灣海平面上升速度是全球 2 倍",
    "   • 2050 年不積極減碳 → 1,398 平方公里被淹",
    "   • 影響約 120 萬人",
    "",
    "✅ 台灣在 COP26 的參與",
    "   • 「台灣日」活動展現能源轉型成果",
    "   • 與友邦國家討論氣候脆弱族群因應之道",
    "",
    "📋 台灣 2050 淨零路徑",
    "   • 能源轉型、產業轉型、生活轉型、社會轉型",
])

# ========== SLIDE 13: 關鍵字回顧 ==========
s13 = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(s13, "📝 COP26 關鍵字，你記住了嗎？", [
    "🌡️ 全球暖化 — 地球發燒了",
    "🧊 海平面上升 — 北極熊的消失",
    "🌱 淨零碳排 — 2050 的目標",
    "⚡ 能源轉型 — 從化石燃料到再生能源",
    "⚖️ 氣候正義 — 沒有人是局外人",
    "👩 性別不平等 — 氣候危機下的弱勢",
    "🏳️ 氣候難民 — 下一個可能是我們",
    "✊ Blah Blah Blah — 年輕世代的怒吼",
], size=24)

# ========== SLIDE 14: 思考題 ==========
s14 = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(s14, "💭 思考 — 寫下你的答案", "2050 年，你會成為怎樣的大人？")
add_body(s14, [
    "🌏 那時的地球是什麼樣子？",
    "",
    "🔑 我們現在做的每一個選擇，都在決定未來",
    "",
    "✍️ 小組討論：",
    "  1. 日常生活中，你可以做哪些事來減碳？",
    "  2. 你認為台灣該如何因應氣候變遷？",
    "  3. 如果有一天你成為總統，你會怎麼做？",
], y=Inches(3.5), size=24)

# ========== SLIDE 15: 結語 ==========
s15 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s15, ACCENT, BG_TOP)
add_title_text(s15, "「改變，從你我開始」", size=44, y=Inches(2.2))
add_subtitle(s15, "永續不是口號，而是每一天的選擇", y=Inches(3.8), size=24)
add_subtitle(s15, "謝謝聆聽 ｜ 讓我們一起行動", y=Inches(5.0), size=20)

output_path = os.path.expanduser(r"~\Desktop\COP26氣候峰會與我們想要的未來_永續素養課.pptx")
prs.save(output_path)
print(f"PPT saved to: {output_path}")
