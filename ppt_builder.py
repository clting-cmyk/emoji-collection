import sys, os, glob, re
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

C_DARK = RGBColor(0x0a, 0x2a, 0x44)
C_MID = RGBColor(0x0a, 0x3d, 0x62)
C_LIGHT = RGBColor(0x1a, 0x6b, 0x8a)
C_ACCENT = RGBColor(0x4f, 0xcd, 0xa5)
C_GOLD = RGBColor(0xff, 0xd7, 0x00)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_OFF = RGBColor(0xE8, 0xF4, 0xF8)
C_ORANGE = RGBColor(0xFF, 0x8c, 0x42)
C_RED = RGBColor(0xef, 0x53, 0x50)
C_PURPLE = RGBColor(0xab, 0x47, 0xbc)

img_dir = os.path.expanduser(r"~\Documents\emoji-collection\ppt_imgs")

def get_img(name):
    import io
    svg_path = None
    for f in os.listdir(img_dir):
        if f.startswith(name) and f.endswith(".svg"):
            svg_path = os.path.join(img_dir, f)
            break
    if not svg_path:
        return None
    png_path = svg_path + ".png"
    if os.path.exists(png_path):
        return png_path
    # Try to convert SVG to PNG using Pillow with SVG support
    try:
        img = Image.open(svg_path)
        img = img.convert("RGBA")
        img.save(png_path, "PNG")
        return png_path
    except Exception:
        pass
    return None

def add_bg(slide, top=C_MID, bottom=C_LIGHT):
    w, h = prs.slide_width, prs.slide_height
    for i in range(15):
        r = top[0] + int((bottom[0] - top[0]) * i / 14)
        g = top[1] + int((bottom[1] - top[1]) * i / 14)
        b = top[2] + int((bottom[2] - top[2]) * i / 14)
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, int(h * i / 15), w, int(h / 15) + 1)
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(r, g, b)
        rect.line.fill.background()

def add_box(slide, l, t, w, h, color, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def add_txt(slide, l, t, w, h, text, size=20, color=C_WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tb

def add_lines(slide, l, t, w, h, lines, size=20, color=C_OFF, gap=Pt(8)):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = gap
    return tb

def add_title_bar(slide, title, subtitle=None):
    add_box(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.2), C_ACCENT)
    add_txt(slide, Inches(0.8), Inches(0.15), Inches(11), Inches(0.8), title, size=36, color=C_DARK, bold=True)
    if subtitle:
        add_txt(slide, Inches(0.8), Inches(0.75), Inches(11), Inches(0.4), subtitle, size=16, color=C_DARK)

def add_img(slide, name, left, top, width, height):
    try:
        img = get_img(name)
        if img:
            slide.shapes.add_picture(img, left, top, width, height)
    except Exception as e:
        print(f"Image error {name}: {e}")

def add_table(slide, rows, cols, left, top, width, height, data):
    ts = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = ts.table
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = data[r][c] if r < len(data) and c < len(data[r]) else ""
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(16)
                p.font.color.rgb = C_WHITE
                p.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_ACCENT if r == 0 else (C_MID if r % 2 == 0 else C_LIGHT)
    return tbl

def add_footer(slide, num, total=15):
    add_txt(slide, Inches(12), Inches(7.0), Inches(1), Inches(0.4), f"{num}/{total}", size=12, color=C_OFF)

# ===================== S1: Cover =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, C_DARK, C_MID)
add_img(s, "earth", Inches(4.8), Inches(0.3), Inches(3.5), Inches(3.5))
add_box(s, Inches(0), Inches(4.0), prs.slide_width, Inches(0.08), C_GOLD)
add_txt(s, Inches(1), Inches(4.3), Inches(11.3), Inches(0.6), "給年輕世代的永續素養課", size=22, color=C_ACCENT, align=PP_ALIGN.CENTER)
add_txt(s, Inches(1), Inches(5.0), Inches(11.3), Inches(1.2), "COP26 氣候峰會\n與我們想要的未來", size=48, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_txt(s, Inches(2), Inches(6.3), Inches(9.3), Inches(0.5), "永續發展 × 氣候行動 × 青年參與", size=20, color=C_GOLD, align=PP_ALIGN.CENTER)
add_txt(s, Inches(4), Inches(6.8), Inches(5.3), Inches(0.4), "對象：高三學生", size=16, color=C_OFF, align=PP_ALIGN.CENTER)
add_footer(s, 1)

# ===================== S2: Flowchart =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title_bar(s, "課程架構", "本課程將帶你一步一步探索氣候變遷")
steps = [
    ("\U0001f30a", "融冰體驗", "感受暖化威脅", C_GOLD),
    ("\U0001f321\ufe0f", "全球暖化", "為什麼發燒？", C_ORANGE),
    ("\U0001f3db\ufe0f", "COP26", "全球氣候峰會", C_ACCENT),
    ("\u2696\ufe0f", "氣候正義", "誰受到影響？", C_PURPLE),
    ("\U0001f331", "青年行動", "你的選擇！", C_GOLD),
]
for i, (emoji, title, desc, color) in enumerate(steps):
    x = Inches(1.0 + i * 2.4)
    y = Inches(2.5)
    add_box(s, x, y, Inches(2.0), Inches(2.5), color, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_txt(s, x, y + Inches(0.2), Inches(2.0), Inches(0.5), emoji, size=36, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_txt(s, x, y + Inches(0.7), Inches(2.0), Inches(0.5), title, size=22, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_txt(s, x, y + Inches(1.3), Inches(2.0), Inches(0.8), desc, size=14, color=C_OFF, align=PP_ALIGN.CENTER)
    if i < 4:
        add_box(s, x + Inches(2.0), y + Inches(1.0), Inches(0.4), Inches(0.3), C_WHITE, MSO_SHAPE.RIGHT_ARROW)
add_img(s, "earth", Inches(9.5), Inches(5.5), Inches(3.5), Inches(1.8))
add_footer(s, 2)

# ===================== S3: Ice melting =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title_bar(s, "\U0001f30a 融冰體驗遊戲", "全球暖化 × 北極熊的家園")
add_lines(s, Inches(0.8), Inches(1.6), Inches(7), Inches(4), [
    "規則（4人一組，站在報紙上）：",
    "",
    "  \u25b8 升溫 0.5\u00b0C → 報紙對折（冰層融化）",
    "  \u25b8 升溫 1.0\u00b0C → 再對折",
    "  \u25b8 升溫 1.5\u00b0C → 再對折",
    "  \u25b8 升溫 2.0\u00b0C → 再對折（誰還站得住？）",
    "",
    "反思：如果你是北極熊，你希望人類怎麼做？",
], size=20)
add_box(s, Inches(8.0), Inches(1.6), Inches(0.04), Inches(4.5), C_ACCENT)
add_lines(s, Inches(8.5), Inches(1.8), Inches(4.5), Inches(3), [
    "科學原理：",
    "",
    "燃燒化石燃料",
    "  \u2193",
    "CO\u2082 增加",
    "  \u2193",
    "溫室效應",
    "  \u2193",
    "地球升溫",
    "  \u2193",
    "冰層融化、海平面上升",
], size=16, color=C_GOLD)
add_img(s, "iceberg", Inches(9.0), Inches(5.0), Inches(4.0), Inches(2.5))
add_footer(s, 3)

# ===================== S4: Global warming =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title_bar(s, "\U0001f321\ufe0f 地球為什麼發燒？", "人類活動造成全球暖化")
add_table(s, 4, 3, Inches(0.8), Inches(1.8), Inches(7.5), Inches(3), [
    ["因素", "升溫幅度", "影響"],
    ["自然因素", "~0.1\u00b0C", "火山活動、太陽輻射變化"],
    ["人為因素", "~1.2\u00b0C", "化石燃料、森林砍伐"],
    ["自然+人為", "~1.3\u00b0C", "目前實際觀測值"],
])
add_img(s, "thermometer", Inches(9.0), Inches(2.0), Inches(3.5), Inches(3.5))
add_txt(s, Inches(0.8), Inches(5.3), Inches(11), Inches(0.8), "關鍵密碼：2030 年碳排減半 \u2502 1.5\u00b0C 安全上限 \u2502 2050 年淨零排放", size=20, color=C_GOLD, bold=True, align=PP_ALIGN.CENTER)
add_footer(s, 4)

# ===================== S5: 1.5 table =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title_bar(s, "爲什麼是 1.5\u00b0C？", "升溫 1.5\u00b0C vs 2\u00b0C 的差距")
add_table(s, 6, 3, Inches(0.8), Inches(1.8), Inches(11.5), Inches(4.5), [
    ["面向", "升溫 1.5\u00b0C", "升溫 2\u00b0C（多 0.5\u00b0C 的代價）"],
    ["極端高溫", "14% 人口受影響", "37% 人口受影響（+17 億人）"],
    ["海平面上升", "上升 0.4 公尺", "上升 0.5 公尺（+1000 萬人）"],
    ["珊瑚礁", "70-90% 退化", "99% 死亡"],
    ["糧食安全", "部分作物減產", "主要作物全面減產"],
    ["經濟損失", "成本降低 25%", "損失大幅增加"],
])
add_footer(s, 5)

# ===================== S6: COP26 intro =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title_bar(s, "\U0001f3db\ufe0f COP26 氣候峰會", "拯救地球的最後機會")
add_lines(s, Inches(0.8), Inches(1.5), Inches(11), Inches(2), [
    "COP = Conference of the Parties（締約國會議）",
    "1995 年起每年舉辦，197 個國家參與",
    "COP26（2021，英國格拉斯哥）— 關鍵轉折點",
], size=22)
add_box(s, Inches(0.8), Inches(3.2), Inches(11.5), Inches(0.04), C_ACCENT)
add_lines(s, Inches(0.8), Inches(3.5), Inches(11), Inches(2.5), [
    "COP26 參與規模",
    "",
    "25,000 名代表（200 國）\u250210,000 名警察\u2502100,000 名抗議民眾",
    "",
    "四大目標：控制 1.5\u00b0C \u2502 1000 億氣候資金 \u2502 保護生態 \u2502 全民合作",
], size=20)
add_img(s, "cop26", Inches(9.5), Inches(4.5), Inches(3.5), Inches(2.5))
add_footer(s, 6)

# ===================== S7: Net zero =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title_bar(s, "\U0001f331 什麼是「淨零排放」？", "Net Zero = 排放量 \u2212 移除量 = 0")
add_lines(s, Inches(0.8), Inches(1.5), Inches(11), Inches(4.5), [
    "淨零 \u2260 零排放",
    "",
    "    排放（Emission）：燃燒化石燃料、工業製程",
    "  \u2212 移除（Removal）：森林碳匯、碳捕集技術",
    "  \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500",
    "    = 淨零（Net Zero）",
    "",
    "全球響應進度：",
    "  128+ 國家 \u00b7 235+ 城市 \u00b7 699+ 企業",
    "  台灣也已宣示：2050 淨零轉型",
], size=20)
add_img(s, "netzero", Inches(9.0), Inches(1.5), Inches(4.0), Inches(3.0))
add_footer(s, 7)

# ===================== S8: COP26 goals table =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title_bar(s, "COP26 四大核心目標", "各國領袖齊聚格拉斯哥的承諾")
add_table(s, 5, 3, Inches(0.8), Inches(1.8), Inches(11.5), Inches(3.5), [
    ["目標", "具體內容", "重要性"],
    ["\u2460 控制升溫", "各國提出更積極減碳承諾", "阻止氣候災難"],
    ["\u2461 氣候資金", "已開發國家提供 1000 億美元", "協助弱勢國家轉型"],
    ["\u2462 保護生態", "森林保育、生物多樣性", "守護自然棲息地"],
    ["\u2463 全民合作", "政府 + 企業 + 民間", "沒有人是局外人"],
])
add_footer(s, 8)

# ===================== S9: Energy transition =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title_bar(s, "\u26a1 能源轉型：淨零之戰的關鍵", "從化石燃料到再生能源")
add_box(s, Inches(0.8), Inches(1.8), Inches(3.5), Inches(1.5), C_RED, MSO_SHAPE.ROUNDED_RECTANGLE)
add_txt(s, Inches(0.8), Inches(2.0), Inches(3.5), Inches(0.5), "化石燃料", size=22, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_txt(s, Inches(0.8), Inches(2.5), Inches(3.5), Inches(0.6), "煤、石油、天然氣\n占全球碳排最大宗", size=14, color=C_OFF, align=PP_ALIGN.CENTER)
add_box(s, Inches(4.5), Inches(2.2), Inches(0.8), Inches(0.5), C_ACCENT, MSO_SHAPE.RIGHT_ARROW)
add_box(s, Inches(5.5), Inches(1.8), Inches(3.5), Inches(1.5), C_ACCENT, MSO_SHAPE.ROUNDED_RECTANGLE)
add_txt(s, Inches(5.5), Inches(2.0), Inches(3.5), Inches(0.5), "再生能源", size=22, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_txt(s, Inches(5.5), Inches(2.5), Inches(3.5), Inches(0.6), "太陽能、風力、水力\n100+ 城市已達 70%", size=14, color=C_OFF, align=PP_ALIGN.CENTER)
add_lines(s, Inches(0.8), Inches(3.8), Inches(11), Inches(3), [
    "綠色城市典範：格拉斯哥（Glasgow）",
    "  \u2022 工業革命搖籃 \u2192 第四波綠色工業革命",
    "  \u2022 蘇格蘭目標：2045 年淨零",
    "  \u2022 2021 年 95% 電力來自再生能源",
    "  \u2022 獲選全球綠色城市",
], size=18)
add_img(s, "wind", Inches(9.5), Inches(3.5), Inches(3.5), Inches(2.5))
add_footer(s, 9)

# ===================== S10: Climate justice =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title_bar(s, "\u2696\ufe0f 氣候正義：誰承受最多？", "造成暖化最少的國家，承受最大的苦難")
add_lines(s, Inches(0.8), Inches(1.5), Inches(11), Inches(4.5), [
    "吐瓦魯 — 第一個面臨消失的國家",
    "  \u2022 平均海拔僅 198 公分",
    "  \u2022 水位每年上升 0.5 公分",
    "  \u2022 外交部長站海水中演講，震撼世界",
    "",
    "性別與氣候",
    "  \u2022 氣候災民中 80% 是女性",
    "  \u2022 400 萬女孩因氣候事件失學",
    "",
    "核心問題：不平等！",
    "  富裕國家排放最多，貧窮國家受害最深",
], size=20)
add_img(s, "justice", Inches(9.5), Inches(1.5), Inches(3.5), Inches(3.0))
add_footer(s, 10)

# ===================== S11: Taiwan =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title_bar(s, "\U0001f30f 氣候難民：下一個可能是我們", "台灣海平面上升速度是全球 2 倍")
add_table(s, 4, 3, Inches(0.8), Inches(1.8), Inches(11), Inches(3), [
    ["面向", "數據", "影響"],
    ["海平面上升", "台灣是全球平均 2 倍", "沿海城市首當其衝"],
    ["淹沒面積", "1,398 平方公里", "約 5 個台北市"],
    ["受影響人口", "約 120 萬人", "2050 年可能發生"],
])
add_img(s, "taiwan", Inches(4.5), Inches(5.0), Inches(4.0), Inches(2.0))
add_footer(s, 11)

# ===================== S12: Youth =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title_bar(s, "\u271a 年輕世代站出來", "No More Blah Blah！行動才有改變")
add_lines(s, Inches(0.8), Inches(1.5), Inches(11), Inches(2.5), [
    "Greta Thunberg（格蕾塔\u00b7童貝里）",
    "  \u2022 2018 年：15 歲，瑞典議會前「為氣候罷課」",
    "  \u2022 2019 年：《時代》風雲人物（史上最年輕）",
    "  \u2022 啟發 400 萬人參與全球氣候遊行",
], size=20)
add_box(s, Inches(0.8), Inches(4.2), Inches(8), Inches(0.04), C_GOLD)
add_lines(s, Inches(0.8), Inches(4.5), Inches(11), Inches(2), [
    "\u201cCOP 已變成公關活動，領導人只會說漂亮話\u201d",
    "  \u2014 Greta 在 COP26 的經典演講",
], size=20, color=C_GOLD)
add_img(s, "youth", Inches(9.0), Inches(1.5), Inches(4.0), Inches(3.0))
add_footer(s, 12)

# ===================== S13: Keywords table =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title_bar(s, "COP26 關鍵字總複習", "這些概念你記住了嗎？")
add_table(s, 9, 2, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5), [
    ["關鍵字", "一句話解釋"],
    ["\U0001f321\ufe0f 全球暖化", "地球因溫室氣體增加而升溫"],
    ["\U0001f9ca 海平面上升", "冰層融化導致海水淹沒陸地"],
    ["\U0001f331 淨零碳排", "排放量 = 移除量"],
    ["\u26a1 能源轉型", "從化石燃料轉向再生能源"],
    ["\u2696\ufe0f 氣候正義", "誰排放最多？誰受害最深？"],
    ["性別不平等", "氣候災民中 80% 是女性"],
    ["\U0001f3f3\ufe0f 氣候難民", "因暖化被迫離開家園的人"],
    ["\u271a Blah Blah Blah", "年輕世代對空頭承諾的不滿"],
])
add_footer(s, 13)

# ===================== S14: Discussion =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title_bar(s, "小組討論：2050 年的你", "思考未來的模樣")
add_lines(s, Inches(0.8), Inches(1.6), Inches(8), Inches(4.5), [
    "2050 年時，你幾歲？地球會是什麼樣子？",
    "",
    "討論問題：",
    "  1. 日常生活中，你可以做哪些事來減碳？",
    "  2. 台灣該如何因應氣候變遷？",
    "  3. 如果你成為總統，你會推動什麼政策？",
    "",
    "想一想：",
    "  你現在做的每一個選擇，",
    "  都在決定 2050 年的地球樣貌。",
], size=22)
add_img(s, "future", Inches(9.5), Inches(1.5), Inches(3.5), Inches(3.0))
add_footer(s, 14)

# ===================== S15: End =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, C_ACCENT, C_DARK)
add_txt(s, Inches(1), Inches(2.0), Inches(11.3), Inches(1.2), "\u201c改變，從你我開始\u201d", size=48, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_box(s, Inches(5), Inches(3.3), Inches(3.3), Inches(0.06), C_GOLD)
add_txt(s, Inches(2), Inches(3.6), Inches(9.3), Inches(0.6), "永續不是口號，而是每一天的選擇", size=24, color=C_GOLD, align=PP_ALIGN.CENTER)
add_img(s, "future", Inches(5.0), Inches(4.5), Inches(3.3), Inches(2.5))
add_txt(s, Inches(2), Inches(6.8), Inches(9.3), Inches(0.5), "謝謝聆聽 \u2502 讓我們一起行動", size=20, color=C_OFF, align=PP_ALIGN.CENTER)
add_footer(s, 15)

output = os.path.expanduser(r"~\Documents\emoji-collection\COP26_Climate_PPT.pptx")
prs.save(output)
print(f"PPT saved: {output}")
print(f"Size: {os.path.getsize(output) / 1024:.0f} KB")
